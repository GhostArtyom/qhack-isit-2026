#!/usr/bin/env python3
"""Generate the final presentation PPTX for ISIT 2026 QHack.
   Uses the official template's slide layouts.

   Usage: uv run python scripts/build_slides.py
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ── helpers ──────────────────────────────────────────────

def add_text(slide, left, top, width, height, text, size=Pt(14), bold=False,
             color=None, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return tf

def add_multi_text(slide, left, top, width, height, lines, size=Pt(14),
                   color=None, spacing=Pt(4), font_name='Calibri'):
    """lines = list of (text, size_override, bold) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, sz, bld = line, size, False
        else:
            text, sz, bld = line[0], line[1] if len(line) > 1 else size, line[2] if len(line) > 2 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = sz
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = spacing
        if color:
            p.font.color.rgb = color
    return tf

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.isabs(path):
        path = os.path.join(PROJECT_ROOT, path)
    if width and height:
        return slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    else:
        return slide.shapes.add_picture(path, left, top)

def add_figure_frame(slide, title, left, top, width, height, border_color=RGBColor(0xDD, 0xDD, 0xDD)):
    """Draw a placeholder rectangle where a figure should go (manual insert)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = border_color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    return shape


# ── build presentation ───────────────────────────────────

def build():
    # Use the official template for its slide dimensions and theme
    template = os.path.join(PROJECT_ROOT,
        'Template - Team Number - ISIT 2026 Quantum Hackathon Presentation.pptx')
    prs = Presentation(template)

    # Remove template placeholder slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Standard slide dimensions (from template)
    W = prs.slide_width   # 9144000 EMU = 10 inches
    H = prs.slide_height  # 5143500 EMU = 5.63 inches

    DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
    ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x66, 0x66, 0x66)

    # We'll create slides from scratch using blank layout
    blank_layout = prs.slide_layouts[6]  # blank is usually index 6

    # ═══════════════════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # Dark background bar
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    add_text(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.6),
             'Quantum-Boosted Channel Allocation', size=Pt(36), bold=True,
             color=WHITE)
    add_text(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.5),
             'for Dense Small-Cell Networks', size=Pt(28), bold=False,
             color=RGBColor(0xBB, 0xCC, 0xEE))
    add_text(slide, Inches(0.8), Inches(1.9), Inches(8.4), Inches(0.4),
             'ISIT 2026 Quantum Hackathon  ·  Team: The Entanglers', size=Pt(16),
             color=RGBColor(0x99, 0xAA, 0xCC))

    add_text(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.6),
             'Warm-Start QAOA for Dynamic QUBO Re-optimization\n'
             'One-Hot vs Binary Encoding: a Qubit-vs-Depth Ansatz Design Choice',
             size=Pt(14), color=GRAY)

    # ═══════════════════════════════════════════════════════
    # SLIDE 2 — PROBLEM
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Problem: Channel Allocation in Dense Small Cells', size=Pt(28), bold=True,
             color=DARK_BLUE)

    # Left column: problem description
    add_multi_text(slide, Inches(0.6), Inches(1.0), Inches(5.0), Inches(2.5),
                   [('Channel allocation = weighted graph coloring → QUBO', Pt(18), True),
                    ('', Pt(8), False),
                    ('• Nodes = base stations, edges = interfering pairs', Pt(15), False),
                    ('• Edge weight = interference coupling (path-loss decay)', Pt(15), False),
                    ('• Assign K channels to N cells → minimize co-channel interference', Pt(15), False),
                    ('• NP-hard: brute force = K^N assignments', Pt(15), False),
                    ('', Pt(8), False),
                    ('CRITICAL: It\'s RECURRING — users move, graph shifts,', Pt(16), True),
                    ('the same QUBO family must be re-solved snapshot after snapshot.', Pt(16), True),
                    ],
                   color=BLACK)

    # Right column: interference graph figure
    fig_path = os.path.join(PROJECT_ROOT, 'fig_interference_graph.png')
    if os.path.exists(fig_path):
        add_image(slide, fig_path, Inches(5.8), Inches(0.8), height=Inches(3.8))

    # Bottom: key insight
    add_text(slide, Inches(0.6), Inches(3.9), Inches(8.8), Inches(0.9),
             'Key insight: Adjacent snapshots differ by a small perturbation → '
             'the previous optimum is near the new one. '
             'We can exploit this structure to cut re-optimization cost.',
             size=Pt(14), bold=False, color=DARK_BLUE)

    # ═══════════════════════════════════════════════════════
    # SLIDE 3 — SOLUTION / WARM-START QAOA
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Solution: Warm-Start QAOA for Dynamic QUBOs', size=Pt(28), bold=True,
             color=DARK_BLUE)

    # Two columns: cold vs warm
    # Cold column
    cold_box = slide.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(4.2), Inches(2.6))
    cold_box.fill.solid()
    cold_box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    cold_box.line.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    add_multi_text(slide, Inches(0.8), Inches(1.1), Inches(3.8), Inches(2.4),
                   [('❄  Cold-Start QAOA (Standard)', Pt(16), True),
                    ('', Pt(6), False),
                    ('Each snapshot: start from uniform superposition', Pt(13), False),
                    ('Optimize all γ, β from scratch via COBYLA', Pt(13), False),
                    ('Every iteration = expensive circuit runs + sampling', Pt(13), False),
                    ('Budget-constrained → far from optimum', Pt(13), False),
                    ],
                   color=BLACK)

    # Warm column
    warm_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.0), Inches(4.2), Inches(2.6))
    warm_box.fill.solid()
    warm_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    warm_box.line.color.rgb = RGBColor(0x27, 0xAE, 0x60)
    add_multi_text(slide, Inches(5.4), Inches(1.1), Inches(3.8), Inches(2.4),
                   [('🔥  Warm-Start QAOA (Ours)', Pt(16), True),
                    ('', Pt(6), False),
                    ('1. Relax prior snapshot\'s bitstring → continuous amplitudes', Pt(13), False),
                    ('2. Initialize ansatz state + mixer from this relaxation', Pt(13), False),
                    ('3. Search starts inside the basin of the prior solution', Pt(13), False),
                    ('→ Better approximation per optimizer iteration ✔', Pt(13), True),
                    ],
                   color=BLACK)

    # Bottom
    add_multi_text(slide, Inches(0.6), Inches(3.8), Inches(8.8), Inches(1.0),
                   [('Based on: Egger, Mareček & Woerner — Warm-starting quantum optimization', Pt(11), False),
                    ('Our contribution: Apply warm-start QAOA to the dynamic/recurring QUBO regime '
                     '+ encoding study. We don\'t claim inventing warm-start QAOA.', Pt(11), False),
                    ],
                   color=GRAY)

    # ═══════════════════════════════════════════════════════
    # SLIDE 4 — ENCODING TRADE-OFF
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Ansatz Design: One-Hot vs Binary Encoding', size=Pt(28), bold=True,
             color=DARK_BLUE)

    # One-hot
    add_multi_text(slide, Inches(0.6), Inches(1.0), Inches(4.2), Inches(2.0),
                   [('One-Hot Encoding', Pt(18), True),
                    ('', Pt(4), False),
                    ('x[v,c] = 1 if cell v uses channel c, else 0', Pt(14), False),
                    ('Qubits: N × K', Pt(14), False),
                    ('Constraint: Σ_c x[v,c] = 1 (shallow penalty)', Pt(14), False),
                    ('✓ Simple constraints, shallow circuit', Pt(14), False),
                    ('✗ Qubit count grows linearly with K', Pt(14), False),
                    ],
                   color=BLACK)
    # Binary
    add_multi_text(slide, Inches(5.2), Inches(1.0), Inches(4.2), Inches(2.0),
                   [('Binary Encoding', Pt(18), True),
                    ('', Pt(4), False),
                    ('Channel index encoded in ⌈log₂K⌉ bits per cell', Pt(14), False),
                    ('Qubits: N × ⌈log₂K⌉', Pt(14), False),
                    ('Constraint: valid indices ≤ K-1 (denser penalty)', Pt(14), False),
                    ('✓ Halves qubit count at K≥4', Pt(14), False),
                    ('✗ Denser cost Hamiltonian, deeper circuit', Pt(14), False),
                    ],
                   color=BLACK)

    # Encoding trade-off figure
    fig_path = os.path.join(PROJECT_ROOT, 'fig_encoding_tradeoff.png')
    if os.path.exists(fig_path):
        add_image(slide, fig_path, Inches(1.0), Inches(3.0), height=Inches(2.0))

    # ═══════════════════════════════════════════════════════
    # SLIDE 5 — HEADLINE RESULTS (Dynamics)
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Headline: Warm-Start Outperforms Cold-Start at Fixed Budget',
             size=Pt(26), bold=True, color=DARK_BLUE)

    # Key numbers
    add_multi_text(slide, Inches(0.6), Inches(1.0), Inches(4.0), Inches(1.5),
                   [('N=8 cells, K=3 channels, 5 mobility snapshots', Pt(16), True),
                    ('Optimizer budget: maxiter=25 (tight)', Pt(14), False),
                    ('', Pt(6), False),
                    ('Mean gap to optimum:', Pt(14), False),
                    ('  Cold-start:  0.178  ✗', Pt(18), True),
                    ('  Warm-start:  0.033  ✓  (5.5× better)', Pt(18), True),
                    ('', Pt(6), False),
                    ('Warm-start nearer optimum on 4/5 snapshots', Pt(14), True),
                    ],
                   color=BLACK)

    # Dynamics figure
    fig_path = os.path.join(PROJECT_ROOT, 'fig_mobility_dynamics.png')
    if os.path.exists(fig_path):
        add_image(slide, fig_path, Inches(4.8), Inches(0.8), height=Inches(4.1))

    # ═══════════════════════════════════════════════════════
    # SLIDE 6 — VERIFICATION
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Verification: 5 Checks Passed', size=Pt(28), bold=True,
             color=DARK_BLUE)

    checks = [
        ('#1  QUBO Correctness',
         'Both encodings: QUBO argmin = brute-force optimum (exact match).',
         '✓'),
        ('#2  Warm-Start Advantage',
         'At fixed maxiter=25, warm-start gap 5.5× smaller than cold-start '
         '(0.033 vs 0.178) across N=8 mobility sequence.',
         '✓'),
        ('#3  Encoding Trade-off',
         'Binary halves qubits vs one-hot (e.g. 12 vs 18 for N=6 K=3). '
         'Depth comparison is non-monotonic — depends on (N,K).',
         '✓'),
        ('#4  Cold QAOA Quality',
         'Best of 3 runs: obj=1.8061 (gap 0.0004 vs brute-force 1.8057). '
         'QAOA consistently matches/beats greedy.',
         '✓'),
        ('#5  Scalability',
         'N=10, K=3: QAOA (binary, 20 qubits) beats greedy by −0.41. '
         'Brute force (3^10 ≈ 59k) intractable for larger graphs.',
         '✓'),
    ]

    y = Inches(1.0)
    for title, desc, status in checks:
        add_text(slide, Inches(0.6), y, Inches(0.5), Inches(0.35),
                 status, size=Pt(20), bold=True, color=RGBColor(0x27, 0xAE, 0x60))
        add_text(slide, Inches(1.1), y, Inches(3.0), Inches(0.3),
                 title, size=Pt(15), bold=True, color=DARK_BLUE)
        add_text(slide, Inches(1.1), y + Inches(0.28), Inches(7.5), Inches(0.35),
                 desc, size=Pt(11), color=GRAY)
        y += Inches(0.72)

    # ═══════════════════════════════════════════════════════
    # SLIDE 7 — LIMITATIONS & NEXT STEPS
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Honest Boundaries & Next Steps', size=Pt(28), bold=True,
             color=DARK_BLUE)

    # Limitations
    add_text(slide, Inches(0.6), Inches(1.0), Inches(4.0), Inches(0.4),
             'Current Limitations', size=Pt(18), bold=True, color=DARK_BLUE)
    add_multi_text(slide, Inches(0.6), Inches(1.4), Inches(4.0), Inches(2.8),
                   [('• Simulator-only: no real quantum hardware yet', Pt(13), False),
                    ('• No noise model: QAOA depth bounded by simulator cost (reps=1)', Pt(13), False),
                    ('• ~26 qubit hard cap: Aer SamplerV2 becomes unreliable beyond', Pt(13), False),
                    ('• λ penalty tuned empirically per graph size', Pt(13), False),
                    ('• Warm-start QAOA is Egger et al.\'s work — we apply it, don\'t invent it', Pt(13), False),
                    ],
                   color=BLACK)

    # Next steps
    add_text(slide, Inches(5.2), Inches(1.0), Inches(4.0), Inches(0.4),
             'Next Steps', size=Pt(18), bold=True, color=DARK_BLUE)
    add_multi_text(slide, Inches(5.2), Inches(1.4), Inches(4.0), Inches(2.8),
                   [('→ Scale to 20+ cell graphs on GPU-accelerated Aer', Pt(13), False),
                    ('→ Add noise model (FakeBrooklyn / FakeSherbrooke)', Pt(13), False),
                    ('→ Explore hardware-native encodings (GRSC, parity)', Pt(13), False),
                    ('→ Test on real IBM hardware for 2-3 qubit proof-of-concept', Pt(13), False),
                    ('→ Extend to multi-objective (throughput + fairness)', Pt(13), False),
                    ],
                   color=BLACK)

    # Bottom
    add_text(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.4),
             'We show: warm-start QAOA is a promising approach for dynamic QUBO regimes '
             '— the hardware to run it at scale is forthcoming.',
             size=Pt(13), bold=False, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # SLIDE 8 — TECH STACK / HACK INGREDIENTS
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.5),
             'Hack Ingredients & Architecture', size=Pt(28), bold=True,
             color=DARK_BLUE)

    add_multi_text(slide, Inches(0.6), Inches(1.0), Inches(4.0), Inches(3.5),
                   [('Quantum Stack', Pt(16), True),
                    ('  Qiskit 2.4  — circuits, transpiler, ansatz', Pt(13), False),
                    ('  qiskit-aer 0.17  — statevector simulator', Pt(13), False),
                    ('  qiskit-algorithms  — QAOA, WarmStartQAOAOptimizer', Pt(13), False),
                    ('  qiskit-optimization  — QuadraticProgram, SlsqpOptimizer', Pt(13), False),
                    ('', Pt(6), False),
                    ('Classical Stack', Pt(16), True),
                    ('  NetworkX  — interference graph generation + mobility', Pt(13), False),
                    ('  NumPy / SciPy  — classical baselines, continuous relaxation', Pt(13), False),
                    ('  COBYLA  — outer-loop parameter optimization', Pt(13), False),
                    ('', Pt(6), False),
                    ('Pipeline', Pt(16), True),
                    ('  Interference graph → coloring QUBO →', Pt(13), False),
                    ('  cold/warm QAOA → decode → evaluate', Pt(13), False),
                    ],
                   color=BLACK)

    # Architecture diagram on the right
    add_multi_text(slide, Inches(5.2), Inches(1.0), Inches(4.2), Inches(4.0),
                   [('Data Flow', Pt(16), True),
                    ('', Pt(6), False),
                    ('src/graphs.py', Pt(14), True),
                    ('  Interference graph + mobility perturbation', Pt(12), False),
                    ('      ↓', Pt(12), False),
                    ('src/qubo.py', Pt(14), True),
                    ('  One-hot / binary QUBO encoding', Pt(12), False),
                    ('      ↓', Pt(12), False),
                    ('src/solve_qaoa.py', Pt(14), True),
                    ('  Cold-start / warm-start QAOA solver', Pt(12), False),
                    ('      ↓', Pt(12), False),
                    ('src/baselines.py  ←→  src/dynamics.py  ←→  src/encoding_study.py', Pt(11), True),
                    ('  Brute force,        Cold vs warm         Qubits vs depth', Pt(11), False),
                    ('  greedy               @ fixed budget       vs feasibility', Pt(11), False),
                    ],
                   color=BLACK)

    # ═══════════════════════════════════════════════════════
    # SLIDE 9 — THANK YOU / Q&A
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(2.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    add_text(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8),
             'Thank You — Questions?', size=Pt(40), bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.5),
             'Team: The Entanglers  ·  ISIT 2026 Quantum Hackathon',
             size=Pt(16), color=RGBColor(0x99, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

    # Contact / reference
    add_multi_text(slide, Inches(1.5), Inches(3.2), Inches(7.0), Inches(1.8),
                   [('Key References', Pt(16), True),
                    ('  Egger, Mareček & Woerner — Warm-starting quantum optimization (arXiv:2009.10095)', Pt(12), False),
                    ('  Qiskit WarmStartQAOAOptimizer documentation', Pt(12), False),
                    ('', Pt(8), False),
                    ('Code & Demo', Pt(16), True),
                    ('  Full pipeline: github.com/.../qhack-isit-2026', Pt(12), False),
                    ('  5 standalone check scripts + 4 auto-generated figures', Pt(12), False),
                    ],
                   color=BLACK)

    # ── save ──────────────────────────────────────────────
    out_path = os.path.join(PROJECT_ROOT, 'Final_Presentation_The_Entanglers.pptx')
    prs.save(out_path)
    print(f'Saved: {out_path}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    build()
